import os
import requests
from flask import Flask, request, jsonify
from dotenv import load_dotenv
import fitz  # PyMuPDF
from pptx import Presentation
import io
import logging
import sys
import json
import uuid
from PIL import Image
import base64

# This line is for local testing. On Render, it does nothing.
load_dotenv()

# --- Configuration ---
API_KEY = os.getenv("GEMINI_API_KEY")
# --- UPDATED --- Using the more powerful gemini-1.5-pro model for better reasoning.
GEMINI_API_URL = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-1.5-pro-latest:generateContent?key={API_KEY}"

app = Flask(__name__)

# --- IMPORTANT NOTE ON CACHE ---
# This in-memory cache works for testing, but it is VOLATILE.
# If your Render service restarts or scales, this data will be lost.
# For a production app, consider using a persistent cache like Redis.
document_cache = {}

# Configure logging
logging.basicConfig(level=logging.INFO, format="[%(asctime)s] %(levelname)s: %(message)s")
logger = logging.getLogger(__name__)


# --- Helper Functions for Text Extraction ---
def extract_text_from_pdf(pdf_file_bytes):
    text = ""
    try:
        with fitz.open(stream=pdf_file_bytes, filetype="pdf") as doc:
            for page in doc:
                text += page.get_text()
        return text
    except Exception as e:
        logger.error(f"PyMuPDF extraction error: {e}")
        return None

def extract_text_from_pptx(pptx_file_stream):
    text = ""
    try:
        prs = Presentation(pptx_file_stream)
        text = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
        return text
    except Exception as e:
        logger.error(f"PPTX extraction error: {e}")
        return None

# --- Helper Functions for AI Processing ---
# --- UPDATED --- This function now accepts a 'history' parameter for conversation memory.
def get_ai_response(prompt, history=None, document_text=None, image_bytes=None, mime_type=None):
    """Handles all communication with the Gemini API, including conversation history."""
    try:
        contents = []
        
        # --- UPDATED --- Build the conversation history for the API
        if history:
            for message in history:
                # The role for Gemini must be 'user' or 'model'
                role = "user" if message.get("role") == "user" else "model"
                contents.append({"role": role, "parts": [{"text": message.get("text")}]})

        # The final user prompt is always the last part of the contents
        final_parts = []
        if document_text:
            context_prompt = (
                "You are an AI study assistant. Use the provided document context to answer the user's query.\n\n"
                f"--- DOCUMENT CONTEXT ---\n{document_text}\n\n"
                f"--- USER QUERY ---\n{prompt}"
            )
            final_parts.append({"text": context_prompt})
        elif image_bytes:
            context_prompt = (
                "You are an AI study assistant. Analyze the following image and answer the user's query about it.\n\n"
                f"--- USER QUERY ---\n{prompt}"
            )
            encoded_image = base64.b64encode(image_bytes).decode('utf-8')
            final_parts.append({"text": context_prompt})
            final_parts.append({"inline_data": {"mime_type": mime_type, "data": encoded_image}})
        else:
            final_parts.append({"text": prompt})
        
        # Add the final prompt as a 'user' turn
        contents.append({"role": "user", "parts": final_parts})

        payload = {"contents": contents}
        headers = {"Content-Type": "application/json"}
        
        response = requests.post(GEMINI_API_URL, headers=headers, json=payload)
        response.raise_for_status()
        
        result_json = response.json()
        
        # Check for safety ratings and blocked responses
        if not result_json.get('candidates'):
            return "The AI could not generate a response, possibly due to safety settings or content restrictions."
        
        return result_json['candidates'][0]['content']['parts'][0]['text']

    except requests.exceptions.HTTPError as http_err:
        logger.error(f"HTTP error occurred: {http_err} - {response.text}")
        return f"Sorry, an error occurred with the AI service (HTTP {response.status_code})."
    except Exception as e:
        logger.error(f"AI generation error: {e}")
        return "Sorry, an error occurred while communicating with the AI."


# --- API Endpoints ---
@app.route('/ai-query', methods=['POST'])
def ai_query():
    """A single endpoint to handle file uploads and text/history queries."""
    
    if 'file' in request.files:
        file = request.files['file']
        filename = file.filename
        file_extension = os.path.splitext(filename)[1].lower()
        file_bytes = file.read()
        
        session_id = str(uuid.uuid4())

        image_extensions = ['.png', '.jpg', '.jpeg', '.webp']
        if file_extension in image_extensions:
            try:
                Image.open(io.BytesIO(file_bytes))
                mime_type = f'image/{file_extension[1:]}' if file_extension.startswith('.') else f'image/{file_extension}'
                document_cache[session_id] = {'type': 'image', 'content': file_bytes, 'mime_type': mime_type}
                return jsonify({"session_id": session_id, "message": "Image processed. You can now ask questions about it."})
            except Exception as e:
                logger.error(f"Image validation error: {e}")
                return jsonify({"error": "Uploaded file is not a valid image."}), 400

        text = None
        if file_extension == '.pdf': text = extract_text_from_pdf(file_bytes)
        elif file_extension == '.pptx': text = extract_text_from_pptx(io.BytesIO(file_bytes))
        elif file_extension == '.txt': text = file_bytes.decode('utf-8')
        else: return jsonify({"error": "Unsupported document type"}), 400

        if not text or not text.strip():
            return jsonify({"error": "Failed to extract text from the document."}), 500

        document_cache[session_id] = {'type': 'text', 'content': text}
        return jsonify({"session_id": session_id, "message": "File processed. You can now ask questions about it."})

    data = request.get_json()
    if not data:
        return jsonify({"error": "Invalid JSON payload."}), 400
        
    query = data.get('query')
    session_id = data.get('session_id')
    history = data.get('history') # --- UPDATED --- Get history from the request

    if not query:
        return jsonify({"error": "Query is required."}), 400

    ai_response = ""
    if session_id:
        cached_data = document_cache.get(session_id)
        if not cached_data:
            return jsonify({"error": "Invalid session ID. Please upload the document again."}), 404
        
        content_type = cached_data.get('type')
        content = cached_data.get('content')

        if content_type == 'text':
            ai_response = get_ai_response(query, history=history, document_text=content)
        elif content_type == 'image':
            mime_type = cached_data.get('mime_type')
            ai_response = get_ai_response(query, history=history, image_bytes=content, mime_type=mime_type)
    else:
        ai_response = get_ai_response(query, history=history) # --- UPDATED --- Pass history
    
    return jsonify({"response": ai_response})

# --- Run the App ---
if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=int(os.environ.get("PORT", 5000)))